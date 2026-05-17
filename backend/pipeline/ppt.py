import os
from typing import List, Dict, Any, Tuple
from .base import PipelineBase, PipelineError

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    Presentation = None
    MSO_SHAPE_TYPE = None


class PptPipeline(PipelineBase):
    SOURCE_TYPE = "ppt"
    SOURCE_ID = "05"

    def parse(self, file_path: str) -> List[Dict[str, Any]]:
        if Presentation is None:
            raise PipelineError("python-pptx not installed")

        try:
            prs = Presentation(file_path)
        except Exception as e:
            raise PipelineError(f"文件损坏，python-pptx 解析失败: {e}")

        doc_file = file_path.split("/")[-1].split("\\")[-1]
        slug = self.doc_slug(doc_file)
        cards = []

        # 图片保存目录
        img_dir = os.path.join(self.data_dir, "raw", "ppt_images", slug)
        os.makedirs(img_dir, exist_ok=True)

        for slide_num, slide in enumerate(prs.slides):
            title = self._extract_title(slide)
            text_content = self._extract_text(slide)

            # 提取图片并做视觉理解
            images_meta, image_descriptions = self._process_images(
                slide, slide_num, img_dir, slug
            )

            # 合并文本 + 图片描述，用 LLM 做摘要
            body = self._build_body(title, text_content, image_descriptions, doc_file)

            if not body.strip():
                continue

            card = self.make_card(
                doc_file=doc_file, doc_slug=slug,
                title=title or f"第{slide_num + 1}页",
                level=0,
                path=self.make_path(doc_file, [title or f"第{slide_num + 1}页"]),
                line_start=slide_num, body=body, seq=slide_num,
            )
            card["slide_num"] = slide_num + 1
            if images_meta:
                card["images"] = images_meta
            cards.append(card)

        if not cards:
            raise PipelineError("PPT无有效文本内容")

        return cards

    def _extract_title(self, slide) -> str:
        if slide.shapes.title and slide.shapes.title.text:
            return slide.shapes.title.text.strip()
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                return shape.text.strip()[:80]
        return ""

    def _extract_text(self, slide) -> str:
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    parts.append(" | ".join(cells))
        return "\n".join(parts)

    def _extract_images(self, slide, max_images: int = 5, min_bytes: int = 5000) -> List[Tuple[bytes, str]]:
        """从 slide 中提取图片，返回 [(image_bytes, content_type), ...]。
        过滤掉过小的图片（logo/icon），限制最大数量。
        """
        images = []
        for shape in slide.shapes:
            if MSO_SHAPE_TYPE and hasattr(shape, 'shape_type'):
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        blob = shape.image.blob
                        ct = shape.image.content_type or "image/png"
                        if blob and len(blob) >= min_bytes:
                            images.append((blob, ct))
                except Exception:
                    continue
            elif hasattr(shape, 'image'):
                try:
                    blob = shape.image.blob
                    ct = shape.image.content_type or "image/png"
                    if blob and len(blob) >= min_bytes:
                        images.append((blob, ct))
                except Exception:
                    continue
            if len(images) >= max_images:
                break
        return images

    def _process_images(
        self, slide, slide_num: int, img_dir: str, slug: str
    ) -> Tuple[List[Dict], List[str]]:
        """提取图片、保存到磁盘、调用视觉模型理解。返回 (images_meta, descriptions)"""
        raw_images = self._extract_images(slide)
        if not raw_images:
            return [], []

        images_meta = []
        descriptions = []

        for img_idx, (blob, content_type) in enumerate(raw_images):
            # 确定文件扩展名
            ext = _content_type_to_ext(content_type)
            filename = f"slide_{slide_num + 1}_img_{img_idx}{ext}"
            save_path = os.path.join(img_dir, filename)

            # 保存图片
            try:
                with open(save_path, "wb") as f:
                    f.write(blob)
            except Exception:
                continue

            # 视觉模型理解
            desc = ""
            try:
                from services.vision_service import describe_image
                desc = describe_image(blob, content_type)
            except Exception as e:
                print(f"[PPT] Vision error slide {slide_num + 1} img {img_idx}: {e}")

            rel_path = os.path.relpath(save_path, self.data_dir)
            images_meta.append({
                "path": rel_path,
                "description": desc[:200] if desc else "",
            })
            if desc:
                descriptions.append(desc)

        return images_meta, descriptions

    def _build_body(
        self,
        title: str,
        text_content: str,
        image_descriptions: List[str],
        doc_file: str,
    ) -> str:
        """合并文本 + 图片描述，用 LLM 生成摘要。"""
        # 如果没有图片，直接返回文本
        if not image_descriptions:
            return text_content

        # 有图片时，调用 LLM 做摘要
        try:
            from services.vision_service import summarize_slide
            summary = summarize_slide(title, text_content, image_descriptions, doc_file)
            if summary and summary.strip():
                return summary.strip()
        except Exception as e:
            print(f"[PPT] Summarize error: {e}")

        # 降级：直接拼接
        parts = []
        if text_content.strip():
            parts.append(text_content.strip())
        for desc in image_descriptions:
            parts.append(f"[图片内容] {desc}")
        return "\n\n".join(parts)


def _content_type_to_ext(content_type: str) -> str:
    mapping = {
        "image/png": ".png",
        "image/jpeg": ".jpg",
        "image/gif": ".gif",
        "image/bmp": ".bmp",
        "image/tiff": ".tiff",
        "image/webp": ".webp",
    }
    return mapping.get(content_type, ".png")
