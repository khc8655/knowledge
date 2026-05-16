import tempfile
import os
import sys
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..'))

from pipeline.ppt import PptPipeline


def test_parse_pptx():
    try:
        from pptx import Presentation
    except ImportError:
        return

    with tempfile.TemporaryDirectory() as tmpdir:
        ppt_path = os.path.join(tmpdir, "test.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Product Overview"
        slide.placeholders[1].text = "AE800 is a 4K video conference system."
        prs.save(ppt_path)

        pipeline = PptPipeline(data_dir=tmpdir)
        cards = pipeline.parse(ppt_path)

        assert len(cards) >= 1
        assert cards[0]["source_type"] == "ppt"
        assert "AE800" in cards[0]["body"]


if __name__ == "__main__":
    test_parse_pptx()
    print("PPT PASS")
